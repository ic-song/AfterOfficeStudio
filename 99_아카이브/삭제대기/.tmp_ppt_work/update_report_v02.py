# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


PPTX = Path(__file__).with_name("report_v02.pptx")
FONT = "맑은 고딕"


def set_text(shape, paragraphs, first_size=13, body_size=11):
    tf = shape.text_frame
    tf.clear()
    for idx, text in enumerate(paragraphs):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = text
        para.level = 0
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(first_size if idx == 0 else body_size)
            run.font.bold = idx == 0


prs = Presentation(PPTX)

# 8페이지: 문서 처리 흐름 - Upstage 적용 위치 보강
s8 = prs.slides[7]
set_text(
    s8.shapes[1],
    [
        "Upstage 기술 적용 위치: Parse/OCR → Chunk → Embedding → Retriever → Solar LLM → Groundedness Check",
        "· Document Parse/OCR은 HWP·PDF·이미지 문서의 본문·표·이미지 내 텍스트 추출 단계에 적용한다.",
        "· Upstage Embedding은 chunk_id 단위로 Vector DB에 적재하고, Solar LLM은 검색 근거 기반 답변·요약·보고서 초안 생성에 활용한다.",
        "· Groundedness Check는 답변이 검색 근거와 일치하는지 검증하는 품질 확인 단계로 둔다.",
    ],
    first_size=12,
    body_size=9,
)
set_text(
    s8.shapes[4],
    [
        "1. 원본 파일 저장 및 문서 ID 부여",
        "→ Object Storage에 원본을 저장하고 document_id를 부여한다.",
        "2. Document Parse/OCR 본문 추출",
        "→ 문서 유형별 Loader와 OCR로 본문, 표, 이미지 내 텍스트를 추출한다.",
        "3. 메타데이터·권한 정보 생성",
        "→ 문서 유형, 부서, 보안등급, 권한, 버전, 최신본 여부를 관리한다.",
        "4. Chunking·Embedding·색인 적재",
        "→ chunk_id 단위로 Search Index와 Vector DB에 적재한다.",
        "5. RAG 답변·검증 이력 관리",
        "→ Retriever, Solar LLM, Groundedness Check 결과와 오류 재처리 이력을 남긴다.",
    ],
    first_size=11,
    body_size=8,
)

# 11페이지: 수행 태스크 - 업로드 운영 시나리오 보강
s11 = prs.slides[10]
for sh in s11.shapes:
    if hasattr(sh, "text") and "업로드 방식" in sh.text and "Loader" in sh.text:
        sh.text = "업로드 방식·Loader 검증"
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.bold = True
set_text(
    s11.shapes[17],
    [
        "운영 시나리오 보강: 관리자 일괄 적재 / 사용자 개별 업로드 / 시스템 API·ETL·Crawler 연계로 태스크를 구분한다.",
        "· 관리자: 기준 문서 묶음 등록, 분류·권한 기본값 지정, 대량 처리 결과와 실패 목록 확인.",
        "· 사용자: 개별 파일 업로드, 문서유형·업무구분 입력, 처리 상태와 검색 반영 여부 확인.",
        "· 시스템 연계: 주기 수집, 변경분 감지, 중복·최신본 판단, 오류 알림과 재처리 로그 관리.",
    ],
    first_size=12,
    body_size=9,
)

# 12페이지: 검증 체크리스트 - 검색/RAG/권한 검증 기준 보강
s12 = prs.slides[11]
set_text(
    s12.shapes[2],
    [
        "검증 기준 보강: 검색 방식별 결과, RAG 답변 근거성, 권한·개인정보 통제를 함께 확인한다.",
        "· 검색: 색인검색, 키워드검색, 벡터검색, 하이브리드 검색의 결과 재현성과 누락 여부를 비교한다.",
        "· RAG: 답변 출처, 원문 페이지/문단, Groundedness Check 통과 여부를 검증한다.",
        "· 보안/운영: 권한 없는 문서 제외, 개인정보 마스킹, 오류 재처리, 처리 이력 추적 가능 여부를 확인한다.",
    ],
    first_size=12,
    body_size=9,
)

# 14페이지: 문서 업로드 운영 기준 - 절차와 완료 기준 보강
s14 = prs.slides[13]
set_text(
    s14.shapes[1],
    [
        "업로드 운영 절차: 업로드 → 형식 검증 → 메타데이터 입력/자동추출 → 권한 확인 → 파싱/청킹 → 임베딩/색인 → 검색/RAG 검증 → 완료/오류 재처리",
        "· 일괄·개별·시스템 연계 업로드 모두 동일한 완료 기준, 실패 재처리 기준, 이력 관리 기준을 적용한다.",
        "· 운영 결과는 처리상태, 검색 반영 여부, 오류 원인, 재처리 결과까지 확인 가능해야 한다.",
    ],
    first_size=11,
    body_size=9,
)

table = s14.shapes[3].table
rows = [
    ["구분", "업로드 방식", "운영 기준", "관리 포인트"],
    ["일괄 업로드", "관리자 대량 등록", "초기 구축·정기 이관·대량 재색인 시 사용", "중복 제거, 실패 재처리, 적재 로그"],
    ["개별 업로드", "사용자 직접 등록", "수시 문서·부서별 업무문서 등록", "권한/분류 입력, 파일 형식 검증"],
    ["시스템 연계", "API·ETL·Crawler", "업무시스템·문서함·외부 법령/자료 연계", "연계 주기, 변경분 감지, 오류 알림"],
    ["완료 기준", "공통 완료 확인", "원본 저장·파싱·청킹·임베딩·색인·검색 검증 완료", "최신본, 보안등급, 출처, 감사 이력"],
]
for r_idx, row in enumerate(rows):
    for c_idx, value in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = value
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(8 if r_idx else 9)
                run.font.bold = r_idx == 0

prs.save(PPTX)
print(PPTX)
