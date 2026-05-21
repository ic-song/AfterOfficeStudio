# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("high_priority_supplement.pptx")
FONT = "맑은 고딕"

NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(67, 126, 181)
SKY = RGBColor(226, 238, 248)
PALE = RGBColor(245, 248, 251)
GRAY = RGBColor(90, 90, 90)
DARK = RGBColor(40, 40, 40)
LINE = RGBColor(190, 205, 220)
ACCENT = RGBColor(217, 83, 79)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def font(run, size=10, bold=False, color=DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, text="", size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        font(r, size=size, bold=bold, color=color)
    return box


def header(slide, title, subtitle, no):
    textbox(slide, 0.55, 0.25, 6.0, 0.25, "소방출동 데이터 기반 AI빅데이터 분석시스템 구축", 8, True)
    textbox(slide, 11.0, 0.25, 1.8, 0.25, "보완자료", 8, True, align=PP_ALIGN.RIGHT)
    textbox(slide, 0.8, 0.72, 7.8, 0.45, title, 20, True)
    textbox(slide, 0.82, 1.15, 8.0, 0.3, subtitle, 9, False, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.48), Inches(11.75), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK
    line.line.color.rgb = DARK
    textbox(slide, 6.18, 7.05, 1.0, 0.2, f"- {no} -", 8, False, GRAY, PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title, body, fill=PALE, title_color=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = title
    for r in p.runs:
        font(r, 11, True, title_color)
    for line in body:
        pp = tf.add_paragraph()
        pp.text = line
        pp.space_before = Pt(2)
        for r in pp.runs:
            font(r, 8.5, False, DARK)
    return shp


def evidence(slide, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.35), Inches(11.75), Inches(0.42))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = ACCENT
    shp.line.dash_style = 2
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(5)
    p = tf.paragraphs[0]
    p.text = "근거  " + text
    for r in p.runs:
        font(r, 9.5, True if r.text.startswith("근거") else False, DARK)


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = BLUE
    conn.line.width = Pt(1.4)
    return conn


def flow_box(slide, x, y, w, h, title, body):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = SKY
    shp.line.color.rgb = BLUE
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        font(r, 9.5, True, NAVY)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.alignment = PP_ALIGN.CENTER
    for r in p2.runs:
        font(r, 7.5, False, DARK)
    return shp


# 1. Page 8 supplement
slide = prs.slides.add_slide(prs.slide_layouts[6])
header(slide, "8페이지 보완｜Upstage 기술 적용 맵", "문서 처리 흐름 안에서 솔루션 적용 위치를 분리해 표시", 1)
items = [
    ("원본 저장", "Object Storage\n문서 ID"),
    ("Parse/OCR", "본문·표·이미지\n텍스트 추출"),
    ("Chunking", "문단/페이지 기준\n분할"),
    ("Embedding", "chunk_id 기준\n벡터화"),
    ("Search/Vector DB", "색인 + Vector DB\n동시 적재"),
    ("Retriever", "근거 문단\n검색/재순위"),
    ("Solar LLM", "답변·요약\n초안 생성"),
    ("Groundedness", "근거 일치성\n검증"),
]
x = 0.75
for i, (t, b) in enumerate(items):
    flow_box(slide, x + i * 1.52, 2.15, 1.25, 0.85, t, b)
    if i < len(items) - 1:
        arrow(slide, x + i * 1.52 + 1.25, 2.57, x + (i + 1) * 1.52 - 0.04, 2.57)
card(slide, 0.9, 3.65, 3.55, 1.3, "보고 메시지", ["문서 적재는 파일 저장이 아니라", "검색/RAG가 사용할 수 있는 지식 데이터 전환 과정"], WHITE)
card(slide, 4.9, 3.65, 3.55, 1.3, "품질 확인 지점", ["OCR·파싱 성공률", "청킹 단위 적정성", "답변 근거 일치성"], WHITE)
card(slide, 8.9, 3.65, 3.55, 1.3, "관리 단위", ["document_id와 chunk_id 기준", "원본, 색인, 벡터, 답변 출처 연결"], WHITE)
evidence(slide, "SFR-005 비정형 문서 구조화/전처리, SFR-017 RAG, SFR-018 AI 지능형 검색")

# 2. Page 11 supplement
slide = prs.slides.add_slide(prs.slide_layouts[6])
header(slide, "11페이지 보완｜업로드 운영 시나리오", "수행 태스크를 운영 주체와 업로드 방식 기준으로 구체화", 2)
card(slide, 0.95, 1.95, 3.45, 3.15, "관리자 일괄 적재", ["기존 SOP·지침·보고서 묶음 등록", "문서유형·부서·권한 기본값 지정", "처리 결과와 실패 목록 확인", "필요 시 재처리 실행"], SKY)
card(slide, 4.95, 1.95, 3.45, 3.15, "사용자 개별 업로드", ["업무 담당자가 수시 문서 등록", "업무구분·문서유형 직접 입력", "처리 상태와 검색 반영 여부 확인", "권한 범위 내 원문 열람"], PALE)
card(slide, 8.95, 1.95, 3.45, 3.15, "시스템 연계 업로드", ["API·ETL·Crawler로 정기 수집", "변경분·최신본·중복 여부 판단", "오류 알림 및 재처리 로그 관리", "연계 주기와 감사 이력 관리"], SKY)
evidence(slide, "SFR-004 문서정보 수집, SFR-005 문서 구조화 및 전처리 자동화")

# 3. Page 12 supplement
slide = prs.slides.add_slide(prs.slide_layouts[6])
header(slide, "12페이지 보완｜검색/RAG 검증 기준", "검색 방식별 확인 기준과 RAG 답변 품질 기준을 분리", 3)
table = slide.shapes.add_table(5, 4, Inches(0.85), Inches(1.85), Inches(11.65), Inches(3.2)).table
headers = ["검증 항목", "확인 기준", "판정 기준", "관련 기능"]
rows = [
    ["색인검색", "문서명·제목·메타데이터 검색", "정확 문서가 상위 노출", "Search Index"],
    ["키워드검색", "본문 키워드·동의어 검색", "누락/오탐 최소화", "Full-text Search"],
    ["벡터검색", "의미 유사 문단 검색", "관련 문단 재현성 확인", "Embedding / Vector DB"],
    ["RAG 답변", "답변·출처·원문 문단 표시", "Groundedness 통과 및 권한 필터 적용", "Retriever / Solar LLM"],
]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
for r, row in enumerate(rows, 1):
    for c, value in enumerate(row):
        cell = table.cell(r, c)
        cell.text = value
        if c == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SKY
for row in table.rows:
    for cell in row.cells:
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                font(run, 8.5, row == table.rows[0], WHITE if row == table.rows[0] else DARK)
card(slide, 0.95, 5.35, 3.5, 0.75, "권한 필터", ["권한 없는 문서는 검색·답변 근거에서 제외"], WHITE)
card(slide, 4.9, 5.35, 3.5, 0.75, "개인정보", ["마스킹 필요 항목은 답변과 원문 노출 전 확인"], WHITE)
card(slide, 8.85, 5.35, 3.5, 0.75, "출처 표시", ["문서명, 페이지, 문단, 최신본 여부를 함께 표시"], WHITE)
evidence(slide, "SFR-017 RAG 기능, SFR-018 AI 지능형 검색 기능, 공공문서 권한·출처 관리 기준")

# 4. Page 14 supplement
slide = prs.slides.add_slide(prs.slide_layouts[6])
header(slide, "14페이지 보완｜문서 업로드 운영 절차", "업로드 방식은 달라도 완료 기준은 동일하게 관리", 4)
steps = ["업로드", "형식 검증", "메타데이터", "권한 확인", "파싱/청킹", "임베딩/색인", "검증", "완료/재처리"]
for i, step in enumerate(steps):
    flow_box(slide, 0.85 + i * 1.46, 2.0, 1.18, 0.72, step, "")
    if i < len(steps) - 1:
        arrow(slide, 0.85 + i * 1.46 + 1.18, 2.36, 0.85 + (i + 1) * 1.46 - 0.05, 2.36)
card(slide, 0.95, 3.35, 2.8, 1.45, "일괄 업로드", ["초기 구축·정기 이관", "대량 문서 재색인", "실패 건 재처리"], SKY)
card(slide, 3.95, 3.35, 2.8, 1.45, "개별 업로드", ["부서별 업무 문서 등록", "문서유형·권한 입력", "처리 상태 확인"], PALE)
card(slide, 6.95, 3.35, 2.8, 1.45, "시스템 연계", ["API·ETL·Crawler", "변경분 감지", "연계 오류 알림"], SKY)
card(slide, 9.95, 3.35, 2.25, 1.45, "운영 결과", ["검색 반영 여부", "오류 원인", "감사 이력"], PALE)
evidence(slide, "SFR-004 문서정보 수집, SFR-005 구조화/전처리, 운영 이력·오류 재처리 관리 필요")

# 5. Page 15 supplement
slide = prs.slides.add_slide(prs.slide_layouts[6])
header(slide, "15페이지 대체｜비정형 데이터 적재 완료 기준", "완료 판단을 원본 저장 이후 검색/RAG 검증까지 확장", 5)
cards = [
    ("원본 저장", ["원본 파일 보관", "document_id 부여"]),
    ("메타데이터", ["문서유형·부서", "권한·버전·최신본"]),
    ("파싱/청킹", ["본문·표·이미지 텍스트", "chunk_id 생성"]),
    ("임베딩/색인", ["Vector DB 적재", "검색 색인 생성"]),
    ("검색/RAG", ["검색 결과 재현성", "답변 출처·Groundedness"]),
    ("권한/운영", ["권한 필터·개인정보", "오류·재처리 이력"]),
]
for i, (t, b) in enumerate(cards):
    card(slide, 0.9 + (i % 3) * 4.0, 1.9 + (i // 3) * 1.65, 3.45, 1.15, t, b, SKY if i < 3 else PALE)
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.55), Inches(11.7), Inches(0.55))
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = ACCENT
tf = box.text_frame
tf.clear()
tf.margin_left = Pt(12)
p = tf.paragraphs[0]
p.text = "완료 판정: 원본, 메타데이터, 파싱/청킹, 임베딩/색인, 검색/RAG, 권한·출처 검증이 모두 확인된 상태"
for r in p.runs:
    font(r, 11, True, DARK)
evidence(slide, "SFR-004, SFR-005, SFR-017, SFR-018")

prs.save(OUT)
print(OUT)
