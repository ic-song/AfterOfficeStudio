# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


PPTX = Path(__file__).with_name("report_v02_readable.pptx")
FONT = "맑은 고딕"


def set_note(shape, title, lines, title_size=13, body_size=10):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(54)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)

    p = tf.paragraphs[0]
    p.text = title
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(title_size)
        run.font.bold = True

    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.level = 0
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(body_size)
            run.font.bold = False


prs = Presentation(PPTX)

# 8페이지: 긴 기술 설명 대신 근거 기반 핵심 연결만 표시
s8 = prs.slides[7]
set_note(
    s8.shapes[1],
    "✓ SFR-005의 Document Loader, Text Split, Vector Embedding, Vector DB, Retriever 요구를 적재 흐름으로 연결.",
    [
        "근거: 요구사항정의서 SFR-005, SFR-017, SFR-018 / ISP 목표모델의 검색서비스 처리 흐름",
    ],
    title_size=12,
    body_size=9,
)

# 11페이지: 수행 태스크에 운영 주체만 짧게 보강
s11 = prs.slides[10]
for sh in s11.shapes:
    if hasattr(sh, "text") and "업로드 방식" in sh.text and "Loader" in sh.text:
        sh.text = "업로드 방식·Loader 검증"
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.bold = True
set_note(
    s11.shapes[17],
    "✓ 수행 태스크는 관리자 일괄 적재, 사용자 개별 업로드, 시스템 연계 업로드 기준으로 구분.",
    [
        "근거: SFR-004 문서정보 수집, SFR-005 비정형 문서 구조화/전처리 자동화",
    ],
    title_size=12,
    body_size=9,
)

# 12페이지: 검증 항목은 검색/RAG 품질 기준 중심으로 간결하게
s12 = prs.slides[11]
set_note(
    s12.shapes[2],
    "✓ 검증은 추출 성공 여부뿐 아니라 검색 결과, 답변 출처, 권한 필터, Groundedness까지 확인.",
    [
        "근거: SFR-017 RAG 기능, SFR-018 AI 지능형 검색 기능, 공공문서 권한·출처 관리 기준",
    ],
    title_size=12,
    body_size=9,
)

# 14페이지: 기존 표는 유지하고 완료 기준만 짧게 보강
s14 = prs.slides[13]
set_note(
    s14.shapes[1],
    "✓ 업로드 운영은 일괄·개별·시스템 연계 방식으로 나누되, 완료 기준은 동일하게 관리.",
    [
        "완료 기준: 원본 저장, 메타데이터 등록, 파싱/청킹, 임베딩/색인, 검색/RAG 검증, 오류 재처리 이력",
    ],
    title_size=12,
    body_size=9,
)

prs.save(PPTX)
print(PPTX)
