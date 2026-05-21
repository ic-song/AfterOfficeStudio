# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("ppt_from_tab_html.pptx")
FONT = "맑은 고딕"

INK = RGBColor(34, 34, 34)
MUTED = RGBColor(102, 102, 102)
LINE = RGBColor(212, 217, 223)
NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(79, 145, 200)
SOFT_BLUE = RGBColor(231, 241, 251)
SOFT = RGBColor(247, 249, 251)
RED = RGBColor(211, 53, 47)
ORANGE = RGBColor(239, 125, 50)
WHITE = RGBColor(255, 255, 255)
DARK_RULE = RGBColor(60, 60, 60)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def set_run(run, size=12, bold=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, value="", size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    for run in p.runs:
        set_run(run, size=size, bold=bold, color=color)
    return shape


def add_paragraph(tf, value, size=12, bold=False, color=INK, level=0, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = value
    p.level = level
    p.alignment = align
    for run in p.runs:
        set_run(run, size=size, bold=bold, color=color)
    return p


def header(slide, title, page_no):
    textbox(slide, 0.55, 0.25, 6.0, 0.25, "소방출동 데이터 기반 AI빅데이터 분석시스템 구축", 8, True)
    textbox(slide, 11.0, 0.25, 1.8, 0.25, "현황분석보고서", 8, True, align=PP_ALIGN.RIGHT)
    textbox(slide, 0.82, 0.78, 7.6, 0.45, title, 20, True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(1.28), Inches(11.7), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = DARK_RULE
    rule.line.color.rgb = DARK_RULE
    textbox(slide, 0.82, 7.05, 1.8, 0.22, "전남소방본부", 9, True)
    textbox(slide, 6.18, 7.05, 1.0, 0.22, f"- 보완 {page_no} -", 8, False, MUTED, PP_ALIGN.CENTER)
    textbox(slide, 10.45, 7.02, 2.1, 0.26, "에코아이티 컨소시엄", 14, False, ORANGE, PP_ALIGN.RIGHT)


def flow_card(slide, x, y, w, h, title, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT_BLUE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        set_run(run, 10, True, NAVY)
    for line in body:
        add_paragraph(tf, line, 7.5, False, INK, align=PP_ALIGN.CENTER)
    return shape


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = BLUE
    conn.line.width = Pt(1.4)
    return conn


def panel(slide, x, y, w, h, title, bullets, fill=WHITE, title_size=17, body_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(157, 187, 215)
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    for run in p.runs:
        set_run(run, title_size, True, NAVY)
    for bullet in bullets:
        add_paragraph(tf, "· " + bullet, body_size, False, INK)
    return shape


def bottom_note(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(6.26), Inches(11.7), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.4)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # Simple target icon
    target = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.02), Inches(6.37), Inches(0.44), Inches(0.44))
    target.fill.background()
    target.line.color.rgb = RED
    target.line.width = Pt(5)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.14), Inches(6.49), Inches(0.20), Inches(0.20))
    inner.fill.background()
    inner.line.color.rgb = RED
    inner.line.width = Pt(3)
    textbox(slide, 1.62, 6.42, 10.5, 0.34, text, 13, False, INK)


def mini_title(slide, value):
    textbox(slide, 0.84, 1.55, 10.5, 0.25, value, 10, True, MUTED)


def add_matrix(slide):
    table = slide.shapes.add_table(5, 4, Inches(0.84), Inches(1.82), Inches(11.7), Inches(2.9)).table
    headers = ["검증 항목", "확인 기준", "판정 기준", "관련 기능"]
    rows = [
        ["색인검색", "문서명·제목·메타데이터 검색", "정확 문서가 상위 노출", "Search Index"],
        ["키워드검색", "본문 키워드·동의어 검색", "누락/오탐 최소화", "Full-text Search"],
        ["벡터검색", "의미 유사 문단 검색", "관련 문단 재현성 확인", "Embedding / Vector DB"],
        ["RAG 답변", "답변·출처·원문 문단 표시", "Groundedness 통과 및 권한 필터 적용", "Retriever / Solar LLM"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            if c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT_BLUE
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, 10 if r_idx == 0 else 9.5, r_idx == 0 or cell == row.cells[0], WHITE if r_idx == 0 else INK)
    return table


# Slide 1
slide = prs.slides.add_slide(blank)
header(slide, "문서 처리 흐름 보완안", 1)
mini_title(slide, "SFR-005의 구조화/전처리 요구를 RAG 활용 데이터 흐름으로 재구성")
steps = [
    ("원본 저장", ["Object Storage", "document_id"]),
    ("Parse/OCR", ["본문·표·이미지", "텍스트 추출"]),
    ("Chunking", ["문단·페이지", "검색 단위 분할"]),
    ("Embedding", ["chunk_id 기준", "벡터화"]),
    ("Search/Vector", ["색인·벡터", "동시 적재"]),
    ("Retriever", ["근거 문단", "검색·재순위"]),
    ("Solar LLM", ["답변·요약", "초안 생성"]),
    ("Groundedness", ["근거 일치성", "검증"]),
]
x0, y0, w, h, gap = 0.84, 2.0, 1.33, 1.02, 0.14
for i, (title, body) in enumerate(steps):
    x = x0 + i * (w + gap)
    flow_card(slide, x, y0, w, h, title, body)
    if i < len(steps) - 1:
        arrow(slide, x + w, y0 + h / 2, x + w + gap, y0 + h / 2)
panel(slide, 0.92, 3.75, 3.45, 1.78, "보고 메시지", ["문서 적재는 파일 저장이 아님", "검색/RAG용 지식 데이터 전환 과정"])
panel(slide, 4.95, 3.75, 3.45, 1.78, "품질 확인 지점", ["OCR·파싱 성공률", "청킹 단위 적정성", "답변 근거 일치성"])
panel(slide, 8.98, 3.75, 3.45, 1.78, "관리 단위", ["document_id와 chunk_id 기준", "원본·색인·벡터·출처 연결"])
bottom_note(slide, "✓ SFR-005 비정형 문서 구조화/전처리, SFR-017 RAG, SFR-018 AI 지능형 검색 요구를 하나의 적재 흐름으로 연결.")

# Slide 2
slide = prs.slides.add_slide(blank)
header(slide, "문서 업로드 운영 시나리오", 2)
mini_title(slide, "수행 태스크를 운영 주체와 업로드 방식 기준으로 구체화")
panel(slide, 0.92, 2.05, 3.45, 3.0, "관리자 일괄 적재", ["기존 SOP·지침·보고서 묶음 등록", "문서유형·부서·권한 기본값 지정", "처리 결과와 실패 목록 확인", "실패 건 재처리 및 이력 확인"], SOFT_BLUE)
panel(slide, 4.95, 2.05, 3.45, 3.0, "사용자 개별 업로드", ["업무 담당자가 수시 문서 등록", "업무구분·문서유형·권한 입력", "처리 상태와 검색 반영 여부 확인", "권한 범위 내 원문 열람"], SOFT)
panel(slide, 8.98, 2.05, 3.45, 3.0, "시스템 연계 업로드", ["API·ETL·Crawler로 정기 수집", "변경분·최신본·중복 여부 판단", "오류 알림과 재처리 로그 관리", "연계 주기와 감사 이력 관리"], SOFT_BLUE)
bottom_note(slide, "✓ SFR-004 문서정보 수집, SFR-005 문서 구조화·전처리 자동화를 운영 시나리오로 구체화.")

# Slide 3
slide = prs.slides.add_slide(blank)
header(slide, "검색/RAG 검증 기준", 3)
mini_title(slide, "검색 방식별 검증 기준과 RAG 답변 품질 기준을 분리")
add_matrix(slide)
panel(slide, 0.92, 5.02, 3.45, 0.85, "권한 필터", ["권한 없는 문서는 검색·답변 근거에서 제외"], WHITE, 14, 10)
panel(slide, 4.95, 5.02, 3.45, 0.85, "개인정보", ["마스킹 필요 항목은 답변·원문 노출 전 확인"], WHITE, 14, 10)
panel(slide, 8.98, 5.02, 3.45, 0.85, "출처 표시", ["문서명, 페이지, 문단, 최신본 여부 표시"], WHITE, 14, 10)
bottom_note(slide, "✓ SFR-017 RAG, SFR-018 AI 지능형 검색, 공공문서 권한·출처 관리 기준을 검증 항목으로 반영.")

# Slide 4
slide = prs.slides.add_slide(blank)
header(slide, "문서 업로드 운영 절차", 4)
mini_title(slide, "업로드 방식은 달라도 완료 기준은 동일하게 관리")
steps = [
    ("업로드", ["일괄·개별", "시스템 연계"]),
    ("형식 검증", ["HWP/PDF", "Office/이미지"]),
    ("메타데이터", ["분류·부서", "버전·최신본"]),
    ("권한 확인", ["보안등급", "접근 범위"]),
    ("파싱/청킹", ["본문 추출", "chunk_id"]),
    ("임베딩/색인", ["Vector DB", "Search Index"]),
    ("검증", ["검색/RAG", "권한 필터"]),
    ("완료/재처리", ["오류 원인", "재처리 이력"]),
]
for i, (title, body) in enumerate(steps):
    x = x0 + i * (w + gap)
    flow_card(slide, x, 2.0, w, h, title, body)
    if i < len(steps) - 1:
        arrow(slide, x + w, 2.0 + h / 2, x + w + gap, 2.0 + h / 2)
panel(slide, 0.92, 3.7, 2.72, 1.72, "일괄 업로드", ["초기 구축·정기 이관", "대량 문서 재색인", "실패 건 재처리"], SOFT_BLUE, 14, 10)
panel(slide, 3.84, 3.7, 2.72, 1.72, "개별 업로드", ["부서별 업무 문서 등록", "문서유형·권한 입력", "처리 상태 확인"], SOFT, 14, 10)
panel(slide, 6.76, 3.7, 2.72, 1.72, "시스템 연계", ["API·ETL·Crawler", "변경분 감지", "연계 오류 알림"], SOFT_BLUE, 14, 10)
panel(slide, 9.68, 3.7, 2.72, 1.72, "운영 결과", ["검색 반영 여부", "오류 원인", "감사 이력"], SOFT, 14, 10)
bottom_note(slide, "✓ SFR-004 문서정보 수집, SFR-005 구조화/전처리, 운영 이력·오류 재처리 관리 기준을 연결.")

# Slide 5
slide = prs.slides.add_slide(blank)
header(slide, "비정형 데이터 적재 완료 기준", 5)
mini_title(slide, "완료 판단을 원본 저장 이후 검색/RAG 검증까지 확장")
completion = [
    ("원본 저장", ["원본 파일 보관", "document_id 부여"]),
    ("메타데이터", ["문서유형·부서", "권한·버전·최신본"]),
    ("파싱/청킹", ["본문·표·이미지 텍스트", "chunk_id 생성"]),
    ("임베딩/색인", ["Vector DB 적재", "검색 색인 생성"]),
    ("검색/RAG", ["검색 결과 재현성", "답변 출처·Groundedness"]),
    ("권한/운영", ["권한 필터·개인정보", "오류·재처리 이력"]),
]
for i, (title, bullets) in enumerate(completion):
    x = 0.92 + (i % 3) * 4.03
    y = 2.0 + (i // 3) * 1.7
    panel(slide, x, y, 3.45, 1.18, title, bullets, SOFT_BLUE if i < 3 else SOFT, 15, 10)
bottom_note(slide, "✓ 완료 판정: 원본, 메타데이터, 파싱/청킹, 임베딩/색인, 검색/RAG, 권한·출처 검증이 모두 확인된 상태.")

prs.save(OUT)
print(OUT)
