# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


PPTX = Path(__file__).with_name("from_original_content_readable.pptx")
FONT = "맑은 고딕"

BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(45, 45, 45)
GRAY = RGBColor(90, 90, 90)
LIGHT_BLUE = RGBColor(231, 241, 250)
LIGHT_GRAY = RGBColor(247, 247, 247)
RED = RGBColor(255, 0, 0)


def apply_font(tf, size=10, bold=False, color=DARK):
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def replace_text(shape, text, size=10, bold=False, color=DARK):
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def set_note(shape, title, lines, title_size=12, body_size=9):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(64)
    tf.margin_right = Pt(18)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = DARK

    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(body_size)
            run.font.bold = False
            run.font.color.rgb = GRAY

    shape.line.color.rgb = RED
    shape.line.dash_style = 2
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.72), Inches(0.52), Inches(6.8), Inches(0.4))
    replace_text(box, title, size=18, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(0.97), Inches(12.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK
    line.line.color.rgb = DARK
    return box


def add_top_labels(slide):
    left = slide.shapes.add_textbox(Inches(0.48), Inches(0.2), Inches(4.7), Inches(0.25))
    replace_text(left, "소방출동 데이터 기반 AI빅데이터 분석시스템 구축", size=9, bold=True)
    right = slide.shapes.add_textbox(Inches(11.25), Inches(0.2), Inches(1.55), Inches(0.25))
    replace_text(right, "현황분석보고서", size=9, bold=True)
    right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_footer(slide, page_text="-14-"):
    left = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(1.5), Inches(0.22))
    replace_text(left, "전남소방본부", size=9, bold=True)
    mid = slide.shapes.add_textbox(Inches(6.2), Inches(7.03), Inches(0.8), Inches(0.2))
    replace_text(mid, page_text, size=8, bold=False)
    mid.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    right = slide.shapes.add_textbox(Inches(10.25), Inches(7.0), Inches(2.2), Inches(0.22))
    replace_text(right, "에코아이티 컨소시엄", size=11, bold=False, color=RGBColor(230, 120, 80))
    right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_card(slide, x, y, w, h, title, body, fill=LIGHT_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(210, 220, 230)
    tf = shape.text_frame
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    for run in p.runs:
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = body
    for run in p2.runs:
        run.font.name = FONT
        run.font.size = Pt(9)
        run.font.color.rgb = DARK
    return shape


prs = Presentation(PPTX)

# 8. 문서 처리 흐름
s8 = prs.slides[7]
set_note(
    s8.shapes[1],
    "✓ SFR-005 처리 흐름에 Upstage 적용 지점(Parse/OCR·Embedding·LLM·Groundedness)을 연결.",
    [
        "근거: 요구사항정의서 SFR-005, SFR-017, SFR-018 / ISP 목표모델 검색서비스 흐름",
    ],
)
replace_text(
    s8.shapes[4],
    "1. 원본 파일 저장 및 문서 ID 부여\n"
    "→ 원본 파일을 저장하고 document_id를 부여함\n\n"
    "2. 문서 본문 텍스트 추출\n"
    "→ Document Parse/OCR로 본문, 표, 이미지 내 텍스트를 추출함\n\n"
    "3. 문서 메타데이터 생성\n"
    "→ 문서 유형, 부서, 권한, 버전, 최신본 여부를 생성함\n\n"
    "4. 검색 색인 및 Vector DB 적재\n"
    "→ 청킹 후 Embedding, Search Index, Vector DB에 적재함\n\n"
    "5. RAG 답변 및 검증 이력 관리\n"
    "→ Retriever, Solar LLM, Groundedness 결과와 재처리 이력을 관리함",
    size=9,
)

# 11. 수행 태스크
s11 = prs.slides[10]
for sh in s11.shapes:
    if hasattr(sh, "text") and "업로드 방식" in sh.text and "Loader" in sh.text:
        replace_text(sh, "업로드 방식·\nLoader 검증", size=10, bold=True)
set_note(
    s11.shapes[17],
    "✓ 운영 시나리오를 관리자 일괄 적재, 사용자 개별 업로드, 시스템 연계 업로드로 구분.",
    [
        "근거: SFR-004 문서정보 수집 / SFR-005 문서 구조화·전처리 자동화",
    ],
)

# 12. 검증 체크리스트
s12 = prs.slides[11]
# 원본 표/그룹은 유지하고 하단 요약만 추가/보강.
note12 = s12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.72),
    Inches(6.05),
    Inches(12.0),
    Inches(0.75),
)
set_note(
    note12,
    "✓ 검증 기준은 추출 성공뿐 아니라 검색 결과, 답변 출처, 권한 필터, Groundedness까지 포함.",
    [
        "근거: SFR-017 RAG 기능 / SFR-018 AI 지능형 검색 기능 / 공공문서 권한·출처 관리 기준",
    ],
)

# 14. 문서 업로드 운영 기준 설계안
s14 = prs.slides[13]
set_note(
    s14.shapes[1],
    "✓ 업로드 방식은 일괄·개별·시스템 연계로 구분하고, 완료 기준은 동일하게 관리.",
    [
        "완료 기준: 원본 저장, 메타데이터 등록, 파싱/청킹, 임베딩/색인, 검색/RAG 검증, 오류 재처리 이력",
    ],
)
table = s14.shapes[3].table
rows = [
    ["구분", "업로드 방식", "운영 기준", "관리 포인트"],
    ["일괄 업로드", "기존 문서 묶음 적재", "초기 구축, 정기 이관, 대량 재색인 시 사용", "중복 제거, 실패 재처리, 적재 로그"],
    ["개별 업로드", "사용자 직접 파일 등록", "수시 문서 등록, 부서별 업무 문서 등록", "권한/분류 입력, 파일 형식 검증"],
    ["시스템 연계", "API, ETL, Crawler 연계", "업무시스템·문서함·외부 법령/자료 연계", "연계 주기, 변경분 감지, 오류 알림"],
    ["완료 기준", "공통 검증", "원본 저장·파싱·청킹·임베딩·색인·검색 검증 완료", "최신본, 보안등급, 출처, 감사 이력"],
]
for r_idx, row in enumerate(rows):
    for c_idx, value in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = value
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(8 if r_idx else 9)
                run.font.bold = r_idx == 0
                run.font.color.rgb = DARK

# 15. 빈 페이지를 완료 기준 페이지로 전환
s15 = prs.slides[14]
add_top_labels(s15)
add_title(s15, "비정형 데이터 적재 완료 기준")

intro = s15.shapes.add_textbox(Inches(0.78), Inches(1.25), Inches(11.5), Inches(0.45))
replace_text(
    intro,
    "비정형 데이터 적재는 원본 보관에서 끝나지 않고, 검색·RAG·출처 검증까지 가능한 상태를 완료 기준으로 본다.",
    size=13,
    bold=True,
)

card_w = Inches(3.75)
card_h = Inches(1.05)
xs = [Inches(0.78), Inches(4.72), Inches(8.66)]
ys = [Inches(1.95), Inches(3.25)]
cards = [
    ("1. 원본·메타데이터", "원본 저장, document_id 부여,\n문서유형·부서·버전·최신본 관리"),
    ("2. 파싱·청킹", "Document Loader/OCR 적용,\n본문·표·이미지 텍스트 추출 및 분할"),
    ("3. 임베딩·색인", "chunk_id 기준 Embedding,\nSearch Index와 Vector DB 적재"),
    ("4. 검색·RAG 검증", "키워드/벡터/하이브리드 검색,\n답변 출처와 Groundedness 확인"),
    ("5. 권한·보안", "부서 권한, 보안등급,\n개인정보 마스킹 및 원문 접근 통제"),
    ("6. 운영 이력", "처리 상태, 오류 원인,\n재처리 결과와 감사 이력 관리"),
]
for i, (title, body) in enumerate(cards):
    add_card(s15, xs[i % 3], ys[i // 3], card_w, card_h, title, body, fill=LIGHT_BLUE if i < 3 else LIGHT_GRAY)

note = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.65), Inches(12.0), Inches(0.75))
set_note(
    note,
    "✓ 완료 판정: 원본 저장, 메타데이터, 파싱/청킹, 임베딩/색인, 검색/RAG, 권한·출처 검증이 모두 확인된 상태.",
    ["근거: SFR-004, SFR-005, SFR-017, SFR-018"],
)
add_footer(s15, "-14-")

# 서식 통일: 본문 텍스트의 기본 글씨체를 맑은 고딕으로 정리
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT

prs.save(PPTX)
print(PPTX)
